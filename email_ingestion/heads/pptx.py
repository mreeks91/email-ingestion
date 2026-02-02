"""PPTX processing head."""

from __future__ import annotations

import io

from email_ingestion.heads.base import HeadInput, HeadResult, Artifact


class PptxHead:
    name = "pptx"
    supported_extensions = {"pptx"}

    def process(self, head_input: HeadInput) -> HeadResult:
        if not head_input.attachment_bytes:
            return HeadResult()
        try:
            from pptx import Presentation  # type: ignore
        except Exception as exc:  # pragma: no cover - import guard
            raise RuntimeError("Missing dependency: python-pptx") from exc
        pres = Presentation(io.BytesIO(head_input.attachment_bytes))
        chunks = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
            if slide.has_notes_slide and slide.notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes and notes.text:
                    chunks.append(notes.text)
        text = "\n".join(chunks).strip() or None
        artifacts = [Artifact(artifact_type="text", text=text)]
        return HeadResult(artifacts=artifacts, metrics={"slides": len(pres.slides)})
